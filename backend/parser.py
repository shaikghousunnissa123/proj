import io
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def parse_pdf(file_bytes: bytes) -> list[dict]:
    """Extracts text from a PDF file, returning a list of dictionaries per page."""
    pages = []
    try:
        pdf_file = io.BytesIO(file_bytes)
        reader = PdfReader(pdf_file)
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append({
                "text": text.strip(),
                "location": f"Page {idx + 1}"
            })
    except Exception as e:
        print(f"Error parsing PDF: {e}")
        # Return at least a fallback raw text if something goes wrong
        pages.append({"text": "", "location": "Page 1"})
    return pages

def parse_docx(file_bytes: bytes) -> list[dict]:
    """Extracts text from a DOCX file."""
    sections = []
    try:
        docx_file = io.BytesIO(file_bytes)
        doc = Document(docx_file)
        
        # Word documents are hard to paginate, so we will chunk paragraphs.
        # We can group paragraphs into sections of ~1500 chars to act as pages.
        current_text = []
        current_len = 0
        section_idx = 1
        
        for para in doc.paragraphs:
            if para.text.strip():
                current_text.append(para.text)
                current_len += len(para.text)
                if current_len > 1500:
                    sections.append({
                        "text": "\n".join(current_text),
                        "location": f"Section {section_idx}"
                    })
                    current_text = []
                    current_len = 0
                    section_idx += 1
                    
        # Append remaining paragraphs
        if current_text:
            sections.append({
                "text": "\n".join(current_text),
                "location": f"Section {section_idx}"
            })
            
        # Parse tables as well and append to the end
        table_text = []
        for table in doc.tables:
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_data:
                    table_text.append(" | ".join(row_data))
        if table_text:
            sections.append({
                "text": "Table Data:\n" + "\n".join(table_text),
                "location": "Tables"
            })
            
    except Exception as e:
        print(f"Error parsing DOCX: {e}")
        sections.append({"text": "", "location": "Document"})
    return sections

def parse_pptx(file_bytes: bytes) -> list[dict]:
    """Extracts text slide-by-slide from a PPTX file."""
    slides = []
    try:
        pptx_file = io.BytesIO(file_bytes)
        prs = Presentation(pptx_file)
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slides.append({
                "text": "\n".join(slide_text),
                "location": f"Slide {idx + 1}"
            })
    except Exception as e:
        print(f"Error parsing PPTX: {e}")
        slides.append({"text": "", "location": "Slide 1"})
    return slides

def parse_txt(file_bytes: bytes) -> list[dict]:
    """Decodes plain text files."""
    try:
        text = file_bytes.decode("utf-8", errors="ignore")
    except Exception:
        try:
            text = file_bytes.decode("latin-1", errors="ignore")
        except Exception as e:
            print(f"Error decoding TXT: {e}")
            text = ""
    return [{"text": text.strip(), "location": "Content"}]

def extract_file_content(filename: str, file_bytes: bytes) -> list[dict]:
    """Determines file type by extension and returns structured page-like blocks."""
    ext = filename.lower().split(".")[-1]
    if ext == "pdf":
        return parse_pdf(file_bytes)
    elif ext == "docx":
        return parse_docx(file_bytes)
    elif ext in ("pptx", "ppt"):
        return parse_pptx(file_bytes)
    elif ext in ("txt", "md", "csv", "json"):
        return parse_txt(file_bytes)
    else:
        # Fallback, try to decode as text
        return parse_txt(file_bytes)
